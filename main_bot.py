import os
import logging
import csv
import io
import json
from datetime import datetime
from flask import Flask, render_template_string
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import Application, CommandHandler, CallbackQueryHandler, ContextTypes

# Настройка логирования
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)

# Конфигурация
BOT_TOKEN = os.environ.get("BOT_TOKEN")
ADMIN_ID = os.environ.get("ADMIN_ID", "")  # ID администратора через запятую
PORT = int(os.environ.get("PORT", 5000))

# Получаем список ID администраторов
admin_ids = [int(x.strip()) for x in ADMIN_ID.split(',')] if ADMIN_ID else []

# Обновленные вопросы согласно уточнению
QUESTIONS = [
    "Рабочая группа Минпросвещения разработала Программу просвещения по запросу родителей.",
    "Программа просвещения родителей – это нормативный документ, по которому должны работать все детские сады.",
    "Один из принципов просвещения – приоритет семьи в вопросах воспитания, обучения и развития.",
    "Основной адресат Программы просвещения – педагоги дошкольных образовательных организаций.",
    "Никто, кроме воспитателей, не может просвещать родителей воспитанников.",
    "Программа просвещения родителей – это новый дополнительный пункт в ФОП ДО",
    "Тематика и формы взаимодействия и педагогического просвещения родителей, которые содержит Программа, – примерные. Педагоги могут их самостоятельно преобразовывать."
]

# Хранилище результатов
class ResultsStorage:
    def __init__(self):
        self.results = {i: {"yes": 0, "no": 0} for i in range(len(QUESTIONS))}
        self.user_progress = {}  # Храним прогресс пользователей
        self.user_answers = {}   # Детальные ответы пользователей
        self.user_info = {}      # Информация о пользователях
    
    def add_vote(self, question_id: int, answer: str, user_id: int, username: str = "", first_name: str = ""):
        # Администраторы не могут участвовать в опросе
        if user_id in admin_ids:
            return False
            
        if question_id in self.results and answer in self.results[question_id]:
            # Обновляем общую статистику
            self.results[question_id][answer] += 1
            
            # Сохраняем прогресс пользователя
            if user_id not in self.user_progress:
                self.user_progress[user_id] = {}
                self.user_answers[user_id] = {}
                self.user_info[user_id] = {
                    "username": username,
                    "first_name": first_name,
                    "last_active": datetime.now().isoformat()
                }
            
            self.user_progress[user_id][question_id] = answer
            self.user_answers[user_id][question_id] = {
                "answer": answer,
                "timestamp": datetime.now().isoformat()
            }
            self.user_info[user_id]["last_active"] = datetime.now().isoformat()
            return True
        return False
    
    def get_user_progress(self, user_id: int):
        return self.user_progress.get(user_id, {})
    
    def get_next_question(self, user_id: int):
        # Администраторы не могут участвовать в опросе
        if user_id in admin_ids:
            return None
            
        user_progress = self.get_user_progress(user_id)
        for i in range(len(QUESTIONS)):
            if i not in user_progress:
                return i
        return None  # Все вопросы пройдены
    
    def get_completion_percentage(self, user_id: int):
        # Администраторы не могут участвовать в опросе
        if user_id in admin_ids:
            return 0
            
        user_progress = self.get_user_progress(user_id)
        return (len(user_progress) / len(QUESTIONS)) * 100
    
    def reset_results(self):
        """Сброс всех результатов (только для админа)"""
        self.results = {i: {"yes": 0, "no": 0} for i in range(len(QUESTIONS))}
        self.user_progress = {}
        self.user_answers = {}
    
    def export_to_csv(self):
        """Экспорт результатов в CSV формат для Google Sheets"""
        output = io.StringIO()
        writer = csv.writer(output)
        
        # Заголовок
        writer.writerow(["Question Number", "Question Text", "Yes", "No", "Total", "Yes %", "No %"])
        
        # Данные по вопросам
        for i, question in enumerate(QUESTIONS):
            stats = self.results[i]
            total = stats["yes"] + stats["no"]
            yes_percent = (stats["yes"] / total * 100) if total > 0 else 0
            no_percent = (stats["no"] / total * 100) if total > 0 else 0
            
            writer.writerow([
                f"Q{i+1}",
                question,
                stats["yes"],
                stats["no"],
                total,
                f"{yes_percent:.1f}%",
                f"{no_percent:.1f}%"
            ])
        
        # Пустая строка
        writer.writerow([])
        
        # Статистика по пользователям
        writer.writerow(["User Statistics"])
        writer.writerow(["User ID", "Username", "Name", "Completed Questions", "Completion %", "Last Active"])
        
        for user_id, info in self.user_info.items():
            completed = len(self.user_progress.get(user_id, {}))
            completion_pct = (completed / len(QUESTIONS)) * 100
            
            writer.writerow([
                user_id,
                info.get("username", ""),
                info.get("first_name", ""),
                completed,
                f"{completion_pct:.1f}%",
                info.get("last_active", "")
            ])
        
        return output.getvalue()
    
    def export_to_html_report(self):
        """Создание интерактивного HTML отчета с графиками"""
        html_template = """
        <!DOCTYPE html>
        <html>
        <head>
            <title>Результаты опроса - Практикум для воспитателей</title>
            <meta charset="utf-8">
            <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
            <style>
                body { font-family: Arial, sans-serif; max-width: 1200px; margin: 0 auto; padding: 20px; }
                .header { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 30px; border-radius: 10px; margin-bottom: 30px; }
                .stats-card { background: white; padding: 20px; border-radius: 10px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); margin-bottom: 20px; }
                .question-card { background: #f8f9fa; padding: 15px; border-radius: 8px; margin: 15px 0; border-left: 4px solid #007bff; }
                .chart-container { height: 300px; margin: 20px 0; }
                .summary-grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 15px; margin: 20px 0; }
                .summary-item { background: white; padding: 15px; border-radius: 8px; text-align: center; box-shadow: 0 2px 5px rgba(0,0,0,0.1); }
                .percentage { font-size: 24px; font-weight: bold; color: #007bff; }
            </style>
        </head>
        <body>
            <div class="header">
                <h1>📊 Результаты опроса</h1>
                <p>Практикум для воспитателей - {{ date }}</p>
                <p>Участников: {{ total_participants }} | Ответов: {{ total_answers }}</p>
            </div>

            <div class="summary-grid">
                <div class="summary-item">
                    <div class="percentage">{{ total_participants }}</div>
                    <div>Участников</div>
                </div>
                <div class="summary-item">
                    <div class="percentage">{{ total_answers }}</div>
                    <div>Всего ответов</div>
                </div>
                <div class="summary-item">
                    <div class="percentage">{{ completion_rate }}%</div>
                    <div>Завершили опрос</div>
                </div>
                <div class="summary-item">
                    <div class="percentage">{{ questions_count }}</div>
                    <div>Вопросов</div>
                </div>
            </div>

            <div class="stats-card">
                <h2>📈 Общая статистика по вопросам</h2>
                <div class="chart-container">
                    <canvas id="overallChart"></canvas>
                </div>
            </div>

            {% for i in range(questions_count) %}
            <div class="stats-card">
                <h3>Вопрос {{ i+1 }}</h3>
                <div class="question-card">
                    <p><strong>{{ questions[i] }}</strong></p>
                </div>
                <div class="chart-container">
                    <canvas id="chart{{ i }}"></canvas>
                </div>
                <p><strong>Результаты:</strong> ✅ Да: {{ results[i].yes }} ({{ yes_percents[i] }}%) | ❌ Нет: {{ results[i].no }} ({{ no_percents[i] }}%)</p>
            </div>
            {% endfor %}

            <script>
                // Общая статистика
                const overallCtx = document.getElementById('overallChart').getContext('2d');
                new Chart(overallCtx, {
                    type: 'bar',
                    data: {
                        labels: {{ question_numbers|tojson }},
                        datasets: [
                            {
                                label: '✅ Да',
                                data: {{ yes_data|tojson }},
                                backgroundColor: '#28a745'
                            },
                            {
                                label: '❌ Нет',
                                data: {{ no_data|tojson }},
                                backgroundColor: '#dc3545'
                            }
                        ]
                    },
                    options: {
                        responsive: true,
                        plugins: {
                            title: {
                                display: true,
                                text: 'Распределение ответов по вопросам'
                            }
                        },
                        scales: {
                            x: {
                                title: {
                                    display: true,
                                    text: 'Номер вопроса'
                                }
                            },
                            y: {
                                title: {
                                    display: true,
                                    text: 'Количество ответов'
                                },
                                beginAtZero: true
                            }
                        }
                    }
                });

                // Графики для каждого вопроса
                {% for i in range(questions_count) %}
                const ctx{{ i }} = document.getElementById('chart{{ i }}').getContext('2d');
                new Chart(ctx{{ i }}, {
                    type: 'doughnut',
                    data: {
                        labels: ['✅ Да ({{ yes_percents[i] }}%)', '❌ Нет ({{ no_percents[i] }}%)'],
                        datasets: [{
                            data: [{{ results[i].yes }}, {{ results[i].no }}],
                            backgroundColor: ['#28a745', '#dc3545']
                        }]
                    },
                    options: {
                        responsive: true,
                        plugins: {
                            legend: {
                                position: 'bottom'
                            },
                            title: {
                                display: true,
                                text: 'Вопрос {{ i+1 }}'
                            }
                        }
                    }
                });
                {% endfor %}
            </script>
        </body>
        </html>
        """
        
        total_answers = sum(sum(stats.values()) for stats in self.results.values())
        total_participants = len(self.user_info)
        
        # Считаем процент завершивших опрос
        completed_users = sum(1 for user_id in self.user_info if len(self.user_progress.get(user_id, {})) == len(QUESTIONS))
        completion_rate = (completed_users / total_participants * 100) if total_participants > 0 else 0
        
        # Подготавливаем данные для графиков
        question_numbers = [f"Вопрос {i+1}" for i in range(len(QUESTIONS))]
        yes_data = [self.results[i]["yes"] for i in range(len(QUESTIONS))]
        no_data = [self.results[i]["no"] for i in range(len(QUESTIONS))]
        
        yes_percents = []
        no_percents = []
        for i in range(len(QUESTIONS)):
            total = self.results[i]["yes"] + self.results[i]["no"]
            yes_percent = (self.results[i]["yes"] / total * 100) if total > 0 else 0
            no_percent = (self.results[i]["no"] / total * 100) if total > 0 else 0
            yes_percents.append(f"{yes_percent:.1f}")
            no_percents.append(f"{no_percent:.1f}")
        
        from flask import render_template_string
        return render_template_string(
            html_template,
            date=datetime.now().strftime("%d.%m.%Y %H:%M"),
            total_participants=total_participants,
            total_answers=total_answers,
            completion_rate=f"{completion_rate:.1f}",
            questions_count=len(QUESTIONS),
            questions=QUESTIONS,
            results=self.results,
            question_numbers=question_numbers,
            yes_data=yes_data,
            no_data=no_data,
            yes_percents=yes_percents,
            no_percents=no_percents
        )
    
    def export_to_pptx(self):
        """Создание простой PowerPoint презентации с результатами"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
            
            prs = Presentation()
            
            # Титульный слайд
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Результаты опроса"
            subtitle.text = f"Практикум для воспитателей\n{datetime.now().strftime('%d.%m.%Y')}\nУчастников: {len(self.user_info)}"
            
            # Слайд с общей статистикой
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Общая статистика"
            
            # Добавляем текстовую статистику
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(1)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            total_answers = sum(sum(stats.values()) for stats in self.results.values())
            p = text_frame.paragraphs[0]
            p.text = f"Всего участников: {len(self.user_info)}\nВсего ответов: {total_answers}\nВопросов: {len(QUESTIONS)}"
            p.font.size = Pt(18)
            
            # Слайды для каждого вопроса
            for i, question in enumerate(QUESTIONS):
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = f"Вопрос {i+1}"
                
                # Текст вопроса
                left = Inches(0.5)
                top = Inches(1)
                width = Inches(9)
                height = Inches(1.5)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                
                p = text_frame.paragraphs[0]
                p.text = question
                p.font.size = Pt(14)
                
                # Статистика
                top = Inches(2.5)
                height = Inches(1)
                
                stats_box = slide.shapes.add_textbox(left, top, width, height)
                stats_frame = stats_box.text_frame
                
                stats = self.results[i]
                total = stats["yes"] + stats["no"]
                yes_percent = (stats["yes"] / total * 100) if total > 0 else 0
                no_percent = (stats["no"] / total * 100) if total > 0 else 0
                
                p = stats_frame.paragraphs[0]
                p.text = f"✅ Да: {stats['yes']} ({yes_percent:.1f}%)\n❌ Нет: {stats['no']} ({no_percent:.1f}%)"
                p.font.size = Pt(16)
                
                # Простая круговая диаграмма (текстовая)
                top = Inches(4)
                height = Inches(1)
                
                chart_box = slide.shapes.add_textbox(left, top, width, height)
                chart_frame = chart_box.text_frame
                
                p = chart_frame.paragraphs[0]
                bar_length = 20
                yes_bars = int(stats["yes"] / total * bar_length) if total > 0 else 0
                no_bars = bar_length - yes_bars
                
                p.text = f"График: [{'█' * yes_bars}{'░' * no_bars}]"
                p.font.size = Pt(12)
            
            # Сохраняем в bytes
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_buffer.seek(0)
            return pptx_buffer
            
        except ImportError:
            # Если библиотека pptx не установлена, создаем текстовый файл с инструкцией
            error_text = """
            Для создания PowerPoint презентаций необходимо установить библиотеку python-pptx.
            
            Добавьте в requirements.txt:
            python-pptx==0.6.21
            
            И перезапустите приложение.
            """
            buffer = io.BytesIO(error_text.encode('utf-8'))
            buffer.name = "INSTALL_INSTRUCTIONS.txt"
            return buffer

results_storage = ResultsStorage()

# Flask приложение для Replit
app = Flask(__name__)

@app.route('/')
def home():
    """Статусная страница для проверки работы бота"""
    html = """
    <!DOCTYPE html>
    <html>
    <head>
        <title>Опрос практикума для воспитателей</title>
        <meta charset="utf-8">
        <style>
            body { font-family: Arial, sans-serif; max-width: 1200px; margin: 0 auto; padding: 20px; }
            .status { background: #f0f8ff; padding: 20px; border-radius: 10px; margin-bottom: 20px; }
            .export-buttons { display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 15px; margin: 20px 0; }
            .export-btn { background: white; padding: 15px; border-radius: 8px; text-align: center; box-shadow: 0 2px 5px rgba(0,0,0,0.1); text-decoration: none; color: #333; border: 2px solid #007bff; }
            .export-btn:hover { background: #007bff; color: white; }
        </style>
    </head>
    <body>
        <h1>🤖 Опрос практикума для воспитателей</h1>
        
        <div class="status">
            <p><strong>Статус:</strong> ✅ Активен</p>
            <p><strong>Версия:</strong> Расширенная с экспортом</p>
            <p><strong>Количество вопросов:</strong> {{ questions_count }}</p>
            <p><strong>Участников:</strong> {{ participants }}</p>
            <p><strong>Всего ответов:</strong> {{ total_answers }}</p>
            <p><strong>Для начала опроса:</strong> Перейдите в Telegram и напишите боту команду <code>/start</code></p>
        </div>

        <h2>📤 Экспорт результатов</h2>
        <div class="export-buttons">
            <a href="/export/html" class="export-btn" target="_blank">
                <strong>🌐 HTML Отчет</strong><br>
                Интерактивный отчет с графиками
            </a>
            <a href="/export/csv" class="export-btn" download>
                <strong>📊 Google Sheets</strong><br>
                CSV для импорта в таблицы
            </a>
            <a href="/export/pptx" class="export-btn" download>
                <strong>📈 PowerPoint</strong><br>
                Презентация с результатами
            </a>
        </div>

        <div class="status">
            <h3>👑 Информация для администраторов:</h3>
            <p>Администраторы не участвуют в опросе, а только управляют статистикой.</p>
            <p>Используйте команду <code>/admin</code> в Telegram для управления.</p>
        </div>
    </body>
    </html>
    """
    
    total_answers = sum(sum(stats.values()) for stats in results_storage.results.values())
    
    return render_template_string(html, 
                                questions_count=len(QUESTIONS),
                                participants=len(results_storage.user_info),
                                total_answers=total_answers)

@app.route('/export/html')
def export_html():
    """Экспорт в HTML отчет"""
    html_content = results_storage.export_to_html_report()
    return html_content

@app.route('/export/csv')
def export_csv():
    """Экспорт в CSV"""
    csv_data = results_storage.export_to_csv()
    response = app.response_class(
        response=csv_data,
        status=200,
        mimetype='text/csv',
        headers={'Content-Disposition': f'attachment; filename=survey_results_{datetime.now().strftime("%Y%m%d_%H%M")}.csv'}
    )
    return response

@app.route('/export/pptx')
def export_pptx():
    """Экспорт в PowerPoint"""
    pptx_buffer = results_storage.export_to_pptx()
    
    if hasattr(pptx_buffer, 'name') and pptx_buffer.name == "INSTALL_INSTRUCTIONS.txt":
        return app.response_class(
            response=pptx_buffer.getvalue(),
            status=200,
            mimetype='text/plain',
            headers={'Content-Disposition': 'attachment; filename=install_instructions.txt'}
        )
    else:
        return app.response_class(
            response=pptx_buffer.getvalue(),
            status=200,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            headers={'Content-Disposition': f'attachment; filename=survey_results_{datetime.now().strftime("%Y%m%d_%H%M")}.pptx'}
        )

@app.route('/health')
def health():
    """Эндпоинт для проверки здоровья приложения"""
    total_answers = sum(sum(stats.values()) for stats in results_storage.results.values())
    return {
        "status": "healthy", 
        "questions_count": len(QUESTIONS),
        "participants": len(results_storage.user_info),
        "total_answers": total_answers,
        "admin_ids": admin_ids
    }

# ... (остальной код бота остается без изменений, только обновляем админ-клавиатуру)

def get_admin_keyboard():
    """Клавиатура для админ панели"""
    keyboard = [
        [InlineKeyboardButton("📊 Статистика", callback_data="admin_stats")],
        [InlineKeyboardButton("📥 Выгрузить CSV", callback_data="admin_export")],
        [InlineKeyboardButton("🌐 HTML Отчет", callback_data="admin_html")],
        [InlineKeyboardButton("📈 PowerPoint", callback_data="admin_pptx")],
        [InlineKeyboardButton("🔄 Сбросить результаты", callback_data="admin_reset")],
        [InlineKeyboardButton("❌ Закрыть", callback_data="admin_close")],
    ]
    return InlineKeyboardMarkup(keyboard)

# Добавляем обработчики для новых кнопок экспорта
async def handle_admin_actions(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обрабатывает действия админа"""
    query = update.callback_query
    user_id = update.effective_user.id
    await query.answer()
    
    if not is_admin(user_id):
        await query.edit_message_text("❌ Доступ запрещен.")
        return
    
    action = query.data
    
    if action == "admin_stats":
        # Показываем детальную статистику
        stats_text = "📊 <b>Детальная статистика:</b>\n\n"
        
        for i in range(len(QUESTIONS)):
            stats = results_storage.results[i]
            total = stats["yes"] + stats["no"]
            yes_percent = (stats["yes"] / total * 100) if total > 0 else 0
            no_percent = (stats["no"] / total * 100) if total > 0 else 0
            
            stats_text += f"<b>Вопрос {i + 1}:</b>\n"
            stats_text += f"✅ Да: {stats['yes']} ({yes_percent:.1f}%)\n"
            stats_text += f"❌ Нет: {stats['no']} ({no_percent:.1f}%)\n"
            stats_text += f"👥 Всего: {total}\n\n"
        
        await query.edit_message_text(stats_text, parse_mode='HTML')
    
    elif action == "admin_export":
        # Выгрузка в CSV
        csv_data = results_storage.export_to_csv()
        csv_file = io.BytesIO(csv_data.encode('utf-8'))
        csv_file.name = f"survey_results_{datetime.now().strftime('%Y%m%d_%H%M')}.csv"
        
        await context.bot.send_document(
            chat_id=user_id,
            document=csv_file,
            filename=csv_file.name,
            caption="📥 <b>Результаты опроса в CSV формате</b>\n\nИмпортируйте в Google Sheets для анализа.",
            parse_mode='HTML'
        )
        
        # Возвращаемся к админ панели
        await admin_command(update, context)
    
    elif action == "admin_html":
        # Создаем HTML отчет и отправляем как файл
        html_content = results_storage.export_to_html_report()
        html_file = io.BytesIO(html_content.encode('utf-8'))
        html_file.name = f"survey_report_{datetime.now().strftime('%Y%m%d_%H%M')}.html"
        
        await context.bot.send_document(
            chat_id=user_id,
            document=html_file,
            filename=html_file.name,
            caption="🌐 <b>Интерактивный HTML отчет</b>\n\nОткройте в браузере для просмотра графиков.",
            parse_mode='HTML'
        )
        
        # Также отправляем ссылку на веб-версию
        web_url = f"https://{os.environ.get('REPL_SLUG', 'your-repl')}.repl.co/export/html"
        await context.bot.send_message(
            chat_id=user_id,
            text=f"🔗 <b>Веб-версия отчета:</b>\n{web_url}",
            parse_mode='HTML'
        )
    
    elif action == "admin_pptx":
        # Выгрузка в PowerPoint
        pptx_buffer = results_storage.export_to_pptx()
        
        if hasattr(pptx_buffer, 'name') and pptx_buffer.name == "INSTALL_INSTRUCTIONS.txt":
            await context.bot.send_document(
                chat_id=user_id,
                document=pptx_buffer,
                filename="install_instructions.txt",
                caption="❌ <b>Библиотека не установлена</b>\n\nСледуйте инструкциям в файле.",
                parse_mode='HTML'
            )
        else:
            pptx_buffer.name = f"survey_results_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
            await context.bot.send_document(
                chat_id=user_id,
                document=pptx_buffer,
                filename=pptx_buffer.name,
                caption="📈 <b>Презентация PowerPoint</b>\n\nГотовая презентация с результатами опроса.",
                parse_mode='HTML'
            )
    
    elif action == "admin_reset":
        # Подтверждение сброса
        confirm_keyboard = InlineKeyboardMarkup([
            [InlineKeyboardButton("✅ Да, сбросить", callback_data="admin_confirm_reset")],
            [InlineKeyboardButton("❌ Отмена", callback_data="admin_cancel_reset")]
        ])
        
        await query.edit_message_text(
            "⚠️ <b>Внимание!</b>\n\n"
            "Вы уверены, что хотите сбросить ВСЕ результаты опроса?\n"
            "Это действие нельзя отменить!",
            reply_markup=confirm_keyboard,
            parse_mode='HTML'
        )
    
    elif action == "admin_confirm_reset":
        # Сброс результатов
        results_storage.reset_results()
        await query.edit_message_text(
            "✅ <b>Все результаты были сброшены!</b>",
            parse_mode='HTML'
        )
    
    elif action == "admin_cancel_reset":
        # Отмена сброса
        await admin_command(update, context)
    
    elif action == "admin_close":
        # Закрытие админ панели
        await query.edit_message_text("👑 Панель администратора закрыта.")

# ... (остальной код бота без изменений)

def main():
    """Основная функция запуска"""
    if not BOT_TOKEN:
        logging.error("BOT_TOKEN не задан в переменных окружения!")
        return
    
    # Создаем приложение бота
    application = Application.builder().token(BOT_TOKEN).build()
    
    # Регистрируем обработчики
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("admin", admin_command))
    application.add_handler(CommandHandler("progress", progress_command))
    application.add_handler(CallbackQueryHandler(handle_answer, pattern="^q[0-9]_(yes|no)$"))
    application.add_handler(CallbackQueryHandler(handle_admin_actions, pattern="^admin_"))
    application.add_handler(CallbackQueryHandler(handle_continue, pattern="^continue_[0-9]$"))
    application.add_error_handler(error_handler)
    
    # Запускаем бота
    logging.info("Бот запускается...")
    logging.info(f"Администраторы: {admin_ids}")
    application.run_polling(drop_pending_updates=True)

if __name__ == "__main__":
    # Запускаем Flask в отдельном потоке для Replit
    from threading import Thread
    flask_thread = Thread(target=lambda: app.run(host='0.0.0.0', port=PORT, debug=False, use_reloader=False))
    flask_thread.daemon = True
    flask_thread.start()
    
    # Запускаем бота в основном потоке
    main()
